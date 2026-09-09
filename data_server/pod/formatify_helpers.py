import json
import mmap
import os
import re
import shutil
import subprocess
import sys
from pathlib import Path
from typing import Dict, List, Optional, Tuple

import mammoth
import pandas as pd
from markdownify import markdownify as md
from pptx import Presentation

from data_server.pod.pod_logger import log_task_error, log_task_info


def _read_csv(file_path: str) -> pd.DataFrame:
    """Read common CSV encodings and let pandas infer the delimiter."""
    last_error = None
    for encoding in ("utf-8-sig", "utf-8", "gb18030"):
        try:
            return pd.read_csv(file_path, encoding=encoding, sep=None, engine="python")
        except UnicodeDecodeError as error:
            last_error = error
    if last_error is not None:
        raise last_error
    return pd.read_csv(file_path, sep=None, engine="python")


def _read_csv_chunked(file_path: str, chunk_size: int):
    """Read CSV in chunks with proper encoding detection."""
    last_error = None
    for encoding in ("utf-8-sig", "utf-8", "gb18030"):
        try:
            return pd.read_csv(
                file_path, 
                encoding=encoding, 
                sep=None, 
                engine="python",
                chunksize=chunk_size,
                iterator=True
            )
        except UnicodeDecodeError as error:
            last_error = error
    if last_error is not None:
        raise last_error
    return pd.read_csv(file_path, sep=None, engine="python", chunksize=chunk_size, iterator=True)


def _count_csv_rows_fast(file_path: str) -> int:
    """Fast CSV row counting without loading into memory."""
    try:
        # Method 1: Use memory mapping (fastest for large files)
        with open(file_path, 'r+b') as f:
            mmapped = mmap.mmap(f.fileno(), 0, access=mmap.ACCESS_READ)
            count = 0
            while mmapped.readline():
                count += 1
            mmapped.close()
        return count - 1  # Exclude header
    except Exception:
        # Method 2: Fallback to simple line counting
        count = 0
        with open(file_path, 'rb') as f:
            for _ in f:
                count += 1
        return count - 1  # Exclude header


def _non_conflicting_output_path(file_path: str) -> str:
    """Avoid overwriting an existing source/target file in the raw stage directory."""
    if not os.path.exists(file_path):
        return file_path
    base_name, extension = os.path.splitext(file_path)
    candidate = f"{base_name}_converted{extension}"
    counter = 1
    while os.path.exists(candidate):
        candidate = f"{base_name}_converted_{counter}{extension}"
        counter += 1
    return candidate


def convert_excel_to_csv(file_path: str, task_uid, use_streaming: bool = False, chunk_size: int = 50000) -> Optional[Dict[str, str]]:
    """
    Convert Excel/CSV to CSV format.
    
    Args:
        file_path: Input file path
        task_uid: Task identifier for logging
        use_streaming: Use streaming mode for large files (default: False)
        chunk_size: Number of rows per chunk in streaming mode (default: 50000)
    
    Returns:
        Conversion result dictionary
    """
    if use_streaming:
        return _convert_excel_to_csv_streaming(file_path, task_uid, chunk_size)
    else:
        return _convert_excel_to_csv_legacy(file_path, task_uid)


def _convert_excel_to_csv_legacy(file_path: str, task_uid) -> Optional[Dict[str, str]]:
    """Original non-streaming implementation."""
    if file_path.lower().endswith(".csv"):
        log_task_info(task_uid, f"CSV source file will be copied without conversion: {file_path}")
        return {
            "from": file_path,
            "to": file_path,
            "to_files": [file_path],
            "status": "success",
        }
    if file_path.lower().endswith((".xlsx", ".xls")):
        log_task_info(task_uid, f"Source file address：{file_path}")
        try:
            # Use openpyxl to avoid file locking issues on Windows
            from openpyxl import load_workbook
            
            # Get sheet names
            wb_info = load_workbook(file_path, read_only=True, data_only=True)
            sheet_names = wb_info.sheetnames
            sheet_count = len(sheet_names)
            wb_info.close()
            
            log_task_info(task_uid, f"Found {sheet_count} sheet(s) in Excel file")
            
            result_files = []
            base_name = os.path.splitext(file_path)[0]
            
            for idx, sheet_name in enumerate(sheet_names, 1):
                try:
                    log_task_info(task_uid, f"Processing sheet {idx}/{sheet_count}: '{sheet_name}'")
                    df = pd.read_excel(file_path, sheet_name=sheet_name)
                    
                    safe_sheet_name = re.sub(r'[<>:"/\\|?*]', '_', sheet_name)
                    if sheet_count == 1:
                        new_file = f"{base_name}.csv"
                    else:
                        new_file = f"{base_name}_{safe_sheet_name}.csv"
                    new_file = _non_conflicting_output_path(new_file)
                    
                    # Use utf-8-sig encoding to ensure Excel can open the CSV correctly
                    df.to_csv(new_file, index=False, encoding='utf-8-sig')
                    result_files.append(new_file)
                    log_task_info(task_uid, f"Sheet '{sheet_name}' converted to {new_file}")
                except Exception as sheet_error:
                    log_task_error(task_uid, f"Failed to convert sheet '{sheet_name}': {sheet_error}")
                    continue
            
            # Remove source file
            try:
                os.remove(file_path)
            except PermissionError:
                # Retry after short delay for Windows file locking
                import time
                time.sleep(0.5)
                try:
                    os.remove(file_path)
                except Exception as e:
                    log_task_error(task_uid, f"Warning: Could not delete source file: {e}")
            
            if len(result_files) == 0:
                return {"from": file_path, "to": None, "status": "failure", "error": "No sheets converted"}
            
            return {
                "from": file_path,
                "to": result_files[0] if len(result_files) == 1 else result_files,
                "to_files": result_files,
                "status": "success",
                "sheets_count": len(result_files)
            }
        except Exception as e:
            log_task_error(task_uid, f"convert file {file_path} error: {e}")
            return {"from": file_path, "to": None, "status": "failure", "error": str(e)}
    return None


def _read_excel_sheet_in_chunks(file_path: str, sheet_name: str, chunk_size: int):
    """
    Read Excel sheet in chunks using openpyxl.
    
    Note: pandas.read_excel() does NOT support chunksize parameter.
    This is a workaround using openpyxl's read_only mode.
    
    Yields:
        DataFrame: Chunk of data as pandas DataFrame
    
    Important:
        This function does NOT load all rows into memory at once.
        It yields chunks as it reads, providing true streaming behavior.
    """
    wb = None
    try:
        from openpyxl import load_workbook
    except ImportError:
        # Fallback: if openpyxl not available, read entire sheet
        df = pd.read_excel(file_path, sheet_name=sheet_name)
        yield df
        return
    
    try:
        # Use read_only mode to save memory
        wb = load_workbook(file_path, read_only=True, data_only=True)
        ws = wb[sheet_name]
        
        # Get header (first row)
        rows_iter = ws.iter_rows(values_only=True)
        header = next(rows_iter, None)
        
        if header is None:
            return
        
        # Convert header to list and clean up
        header = [str(cell) if cell is not None else f"Column_{i}" for i, cell in enumerate(header)]
        
        # Stream processing: yield chunks without loading all rows
        chunk_data = []
        for row in rows_iter:
            chunk_data.append(row)
            
            if len(chunk_data) >= chunk_size:
                # Convert to DataFrame and yield
                df = pd.DataFrame(chunk_data, columns=header)
                yield df
                chunk_data = []
                del df
        
        # Yield remaining data
        if chunk_data:
            df = pd.DataFrame(chunk_data, columns=header)
            yield df
            del df
        
    except Exception as e:
        # Fallback to reading entire sheet
        df = pd.read_excel(file_path, sheet_name=sheet_name)
        yield df
    finally:
        # Ensure workbook is closed
        if wb is not None:
            try:
                wb.close()
            except Exception:
                pass


def _convert_excel_to_csv_streaming(file_path: str, task_uid, chunk_size: int) -> Optional[Dict[str, str]]:
    """
    Streaming implementation with chunked reading and writing.
    
    Note: pandas.read_excel() does NOT support chunksize parameter.
    We use openpyxl's read_only mode as a workaround.
    """
    if file_path.lower().endswith(".csv"):
        log_task_info(task_uid, f"CSV source file will be copied without conversion: {file_path}")
        return {
            "from": file_path,
            "to": file_path,
            "to_files": [file_path],
            "status": "success",
        }
    
    if file_path.lower().endswith((".xlsx", ".xls")):
        log_task_info(task_uid, f"[Streaming Mode] Source file address: {file_path}")
        log_task_info(task_uid, f"[Streaming Mode] Chunk size: {chunk_size:,} rows")
        
        try:
            # Use openpyxl directly to avoid file handle issues
            from openpyxl import load_workbook
            
            # First pass: get sheet names
            wb_info = load_workbook(file_path, read_only=True, data_only=True)
            sheet_names = wb_info.sheetnames
            sheet_count = len(sheet_names)
            wb_info.close()
            
            log_task_info(task_uid, f"Found {sheet_count} sheet(s) in Excel file")
            
            result_files = []
            base_name = os.path.splitext(file_path)[0]
            
            for idx, sheet_name in enumerate(sheet_names, 1):
                new_file = None
                try:
                    log_task_info(task_uid, f"Processing sheet {idx}/{sheet_count}: '{sheet_name}'")
                    
                    # Generate output filename
                    safe_sheet_name = re.sub(r'[<>:"/\\|?*]', '_', sheet_name)
                    if sheet_count == 1:
                        new_file = f"{base_name}.csv"
                    else:
                        new_file = f"{base_name}_{safe_sheet_name}.csv"
                    new_file = _non_conflicting_output_path(new_file)
                    
                    # Count total rows for progress reporting using openpyxl
                    log_task_info(task_uid, f"Counting total rows in sheet '{sheet_name}'...")
                    wb_count = load_workbook(file_path, read_only=True, data_only=True)
                    ws_count = wb_count[sheet_name]
                    total_rows = ws_count.max_row - 1  # Exclude header
                    wb_count.close()
                    log_task_info(task_uid, f"Total rows: {total_rows:,}")
                    
                    # Stream processing using openpyxl
                    first_chunk = True
                    rows_processed = 0
                    
                    for chunk in _read_excel_sheet_in_chunks(file_path, sheet_name, chunk_size):
                        chunk.to_csv(
                            new_file,
                            mode='w' if first_chunk else 'a',
                            header=first_chunk,
                            index=False,
                            encoding='utf-8-sig'
                        )
                        first_chunk = False
                        rows_processed += len(chunk)
                        
                        # Progress reporting
                        progress = (rows_processed / total_rows * 100) if total_rows > 0 else 0
                        log_task_info(
                            task_uid,
                            f"Progress: {rows_processed:,}/{total_rows:,} rows ({progress:.1f}%)"
                        )
                        
                        # Release memory
                        del chunk
                    
                    result_files.append(new_file)
                    log_task_info(task_uid, f"Sheet '{sheet_name}' converted successfully to {new_file}")
                    
                except Exception as sheet_error:
                    log_task_error(task_uid, f"Failed to convert sheet '{sheet_name}': {sheet_error}")
                    # Rollback: delete partial file
                    if new_file and os.path.exists(new_file):
                        try:
                            os.remove(new_file)
                            log_task_info(task_uid, f"Rolled back partial file: {new_file}")
                        except Exception:
                            pass
                    continue
            
            # Remove source file only if at least one sheet succeeded
            if len(result_files) > 0:
                try:
                    os.remove(file_path)
                except PermissionError:
                    # File might be locked, try again after a short delay
                    import time
                    time.sleep(0.5)
                    try:
                        os.remove(file_path)
                    except Exception as e:
                        log_task_error(task_uid, f"Warning: Could not delete source file: {e}")
            
            if len(result_files) == 0:
                return {"from": file_path, "to": None, "status": "failure", "error": "No sheets converted"}
            
            return {
                "from": file_path,
                "to": result_files[0] if len(result_files) == 1 else result_files,
                "to_files": result_files,
                "status": "success",
                "sheets_count": len(result_files)
            }
        except Exception as e:
            log_task_error(task_uid, f"convert file {file_path} error: {e}")
            return {"from": file_path, "to": None, "status": "failure", "error": str(e)}
    return None


def convert_excel_to_json(file_path: str, task_uid, use_streaming: bool = False, chunk_size: int = 50000) -> Optional[Dict[str, str]]:
    """
    Convert Excel/CSV to JSON format.
    
    Args:
        file_path: Input file path
        task_uid: Task identifier for logging
        use_streaming: Use streaming mode for large files (default: False)
        chunk_size: Number of rows per chunk in streaming mode (default: 50000)
    
    Returns:
        Conversion result dictionary
    """
    if use_streaming:
        return _convert_excel_to_json_streaming(file_path, task_uid, chunk_size)
    else:
        return _convert_excel_to_json_legacy(file_path, task_uid)


def _convert_excel_to_json_legacy(file_path: str, task_uid) -> Optional[Dict[str, str]]:
    """Original non-streaming implementation."""
    if file_path.lower().endswith(".csv"):
        log_task_info(task_uid, f"Source file address: {file_path}")
        try:
            new_file = _non_conflicting_output_path(
                f"{os.path.splitext(file_path)[0]}.json"
            )
            _read_csv(file_path).to_json(new_file, orient="records", force_ascii=False)
            return {
                "from": file_path,
                "to": new_file,
                "to_files": [new_file],
                "status": "success",
            }
        except Exception as e:
            log_task_error(task_uid, f"convert file {file_path} error: {e}")
            return {"from": file_path, "to": None, "status": "failure", "error": str(e)}
    if file_path.lower().endswith((".xlsx", ".xls")):
        log_task_info(task_uid, f"Source file address：{file_path}")
        try:
            xls = pd.ExcelFile(file_path)
            sheet_names = xls.sheet_names
            sheet_count = len(sheet_names)
            
            log_task_info(task_uid, f"Found {sheet_count} sheet(s) in Excel file")
            
            result_files = []
            base_name = os.path.splitext(file_path)[0]
            
            for idx, sheet_name in enumerate(sheet_names, 1):
                try:
                    log_task_info(task_uid, f"Processing sheet {idx}/{sheet_count}: '{sheet_name}'")
                    df = pd.read_excel(file_path, sheet_name=sheet_name)
                    
                    safe_sheet_name = re.sub(r'[<>:"/\\|?*]', '_', sheet_name)
                    if sheet_count == 1:
                        new_file = f"{base_name}.json"
                    else:
                        new_file = f"{base_name}_{safe_sheet_name}.json"
                    new_file = _non_conflicting_output_path(new_file)
                    
                    df.to_json(new_file, orient="records", force_ascii=False)
                    result_files.append(new_file)
                    log_task_info(task_uid, f"Sheet '{sheet_name}' converted to {new_file}")
                except Exception as sheet_error:
                    log_task_error(task_uid, f"Failed to convert sheet '{sheet_name}': {sheet_error}")
                    continue
            
            os.remove(file_path)
            
            if len(result_files) == 0:
                return {"from": file_path, "to": None, "status": "failure", "error": "No sheets converted"}
            
            return {
                "from": file_path,
                "to": result_files[0] if len(result_files) == 1 else result_files,
                "to_files": result_files,
                "status": "success",
                "sheets_count": len(result_files)
            }
        except Exception as e:
            log_task_error(task_uid, f"convert file {file_path} error: {e}")
            return {"from": file_path, "to": None, "status": "failure", "error": str(e)}
    return None


def _convert_excel_to_json_streaming(file_path: str, task_uid, chunk_size: int) -> Optional[Dict[str, str]]:
    """
    Streaming implementation with chunked reading and writing.
    
    Note: pandas.read_excel() does NOT support chunksize parameter.
    We use openpyxl's read_only mode as a workaround.
    """
    if file_path.lower().endswith(".csv"):
        log_task_info(task_uid, f"[Streaming Mode] Source file address: {file_path}")
        log_task_info(task_uid, f"[Streaming Mode] Chunk size: {chunk_size:,} rows")
        
        new_file = None
        try:
            new_file = _non_conflicting_output_path(
                f"{os.path.splitext(file_path)[0]}.json"
            )
            
            # Count total rows
            log_task_info(task_uid, "Counting total rows...")
            total_rows = _count_csv_rows_fast(file_path)
            log_task_info(task_uid, f"Total rows: {total_rows:,}")
            
            # Stream processing
            first_chunk = True
            rows_processed = 0
            
            with open(new_file, 'w', encoding='utf-8') as json_file:
                json_file.write('[')
                
                for chunk in _read_csv_chunked(file_path, chunk_size):
                    json_str = chunk.to_json(orient='records', force_ascii=False, indent=None)
                    json_str = json_str[1:-1]  # Remove [ ]
                    
                    if json_str:
                        if not first_chunk:
                            json_file.write(',')
                        json_file.write(json_str)
                        first_chunk = False
                    
                    rows_processed += len(chunk)
                    progress = (rows_processed / total_rows * 100) if total_rows > 0 else 0
                    log_task_info(
                        task_uid,
                        f"Progress: {rows_processed:,}/{total_rows:,} rows ({progress:.1f}%)"
                    )
                    
                    del chunk
                
                json_file.write(']')
            
            return {
                "from": file_path,
                "to": new_file,
                "to_files": [new_file],
                "status": "success",
            }
        except Exception as e:
            log_task_error(task_uid, f"convert file {file_path} error: {e}")
            # Rollback
            if new_file and os.path.exists(new_file):
                try:
                    os.remove(new_file)
                except Exception:
                    pass
            return {"from": file_path, "to": None, "status": "failure", "error": str(e)}
    
    if file_path.lower().endswith((".xlsx", ".xls")):
        log_task_info(task_uid, f"[Streaming Mode] Source file address: {file_path}")
        log_task_info(task_uid, f"[Streaming Mode] Chunk size: {chunk_size:,} rows")
        
        try:
            from openpyxl import load_workbook
            
            # First pass: get sheet names
            wb_info = load_workbook(file_path, read_only=True, data_only=True)
            sheet_names = wb_info.sheetnames
            sheet_count = len(sheet_names)
            wb_info.close()
            
            log_task_info(task_uid, f"Found {sheet_count} sheet(s) in Excel file")
            
            result_files = []
            base_name = os.path.splitext(file_path)[0]
            
            for idx, sheet_name in enumerate(sheet_names, 1):
                new_file = None
                try:
                    log_task_info(task_uid, f"Processing sheet {idx}/{sheet_count}: '{sheet_name}'")
                    
                    safe_sheet_name = re.sub(r'[<>:"/\\|?*]', '_', sheet_name)
                    if sheet_count == 1:
                        new_file = f"{base_name}.json"
                    else:
                        new_file = f"{base_name}_{safe_sheet_name}.json"
                    new_file = _non_conflicting_output_path(new_file)
                    
                    # Count total rows using openpyxl
                    log_task_info(task_uid, f"Counting total rows in sheet '{sheet_name}'...")
                    wb_count = load_workbook(file_path, read_only=True, data_only=True)
                    ws_count = wb_count[sheet_name]
                    total_rows = ws_count.max_row - 1  # Exclude header
                    wb_count.close()
                    log_task_info(task_uid, f"Total rows: {total_rows:,}")
                    
                    # Stream processing
                    first_chunk = True
                    rows_processed = 0
                    
                    with open(new_file, 'w', encoding='utf-8') as json_file:
                        json_file.write('[')
                        
                        for chunk in _read_excel_sheet_in_chunks(file_path, sheet_name, chunk_size):
                            json_str = chunk.to_json(orient='records', force_ascii=False, indent=None)
                            json_str = json_str[1:-1]
                            
                            if json_str:
                                if not first_chunk:
                                    json_file.write(',')
                                json_file.write(json_str)
                                first_chunk = False
                            
                            rows_processed += len(chunk)
                            progress = (rows_processed / total_rows * 100) if total_rows > 0 else 0
                            log_task_info(
                                task_uid,
                                f"Progress: {rows_processed:,}/{total_rows:,} rows ({progress:.1f}%)"
                            )
                            
                            del chunk
                        
                        json_file.write(']')
                    
                    result_files.append(new_file)
                    log_task_info(task_uid, f"Sheet '{sheet_name}' converted successfully to {new_file}")
                    
                except Exception as sheet_error:
                    log_task_error(task_uid, f"Failed to convert sheet '{sheet_name}': {sheet_error}")
                    # Rollback
                    if new_file and os.path.exists(new_file):
                        try:
                            os.remove(new_file)
                            log_task_info(task_uid, f"Rolled back partial file: {new_file}")
                        except Exception:
                            pass
                    continue
            
            if len(result_files) > 0:
                try:
                    os.remove(file_path)
                except PermissionError:
                    import time
                    time.sleep(0.5)
                    try:
                        os.remove(file_path)
                    except Exception as e:
                        log_task_error(task_uid, f"Warning: Could not delete source file: {e}")
            
            if len(result_files) == 0:
                return {"from": file_path, "to": None, "status": "failure", "error": "No sheets converted"}
            
            return {
                "from": file_path,
                "to": result_files[0] if len(result_files) == 1 else result_files,
                "to_files": result_files,
                "status": "success",
                "sheets_count": len(result_files)
            }
        except Exception as e:
            log_task_error(task_uid, f"convert file {file_path} error: {e}")
            return {"from": file_path, "to": None, "status": "failure", "error": str(e)}
    return None


def convert_excel_to_parquet(file_path: str, task_uid, use_streaming: bool = False, chunk_size: int = 50000) -> Optional[Dict[str, str]]:
    """
    Convert Excel/CSV to Parquet format.
    
    Args:
        file_path: Input file path
        task_uid: Task identifier for logging
        use_streaming: Use streaming mode for large files (default: False)
        chunk_size: Number of rows per chunk in streaming mode (default: 50000)
    
    Returns:
        Conversion result dictionary
    """
    if use_streaming:
        return _convert_excel_to_parquet_streaming(file_path, task_uid, chunk_size)
    else:
        return _convert_excel_to_parquet_legacy(file_path, task_uid)


def _convert_excel_to_parquet_legacy(file_path: str, task_uid) -> Optional[Dict[str, str]]:
    """Original non-streaming implementation."""
    if file_path.lower().endswith(".csv"):
        log_task_info(task_uid, f"Source file address: {file_path}")
        try:
            new_file = _non_conflicting_output_path(
                f"{os.path.splitext(file_path)[0]}.parquet"
            )
            df = _read_csv(file_path)
            for col in df.columns:
                if df[col].dtype == "object":
                    df[col] = df[col].astype(str).replace("nan", None)
                elif pd.api.types.is_integer_dtype(df[col]) and df[col].isna().any():
                    df[col] = df[col].astype(str)
            df.to_parquet(new_file, index=False, engine="pyarrow")
            return {
                "from": file_path,
                "to": new_file,
                "to_files": [new_file],
                "status": "success",
            }
        except Exception as e:
            log_task_error(task_uid, f"convert file {file_path} error: {e}")
            return {"from": file_path, "to": None, "status": "failure", "error": str(e)}
    if file_path.lower().endswith((".xlsx", ".xls")):
        log_task_info(task_uid, f"Source file address：{file_path}")
        try:
            # Read Excel file to get all sheet names
            xls = pd.ExcelFile(file_path)
            sheet_names = xls.sheet_names
            sheet_count = len(sheet_names)
            
            log_task_info(task_uid, f"Found {sheet_count} sheet(s) in Excel file")
            
            result_files = []
            base_name = os.path.splitext(file_path)[0]
            
            # Process each sheet
            for idx, sheet_name in enumerate(sheet_names, 1):
                try:
                    log_task_info(task_uid, f"Processing sheet {idx}/{sheet_count}: '{sheet_name}'")
                    
                    # Read the sheet
                    df = pd.read_excel(file_path, sheet_name=sheet_name)
                    
                    # Data type processing (same as original logic)
                    for col in df.columns:
                        if df[col].dtype == "object":
                            df[col] = df[col].astype(str)
                            df[col] = df[col].replace("nan", None)
                        elif pd.api.types.is_integer_dtype(df[col]):
                            if df[col].isna().any():
                                df[col] = df[col].astype(str)
                        elif pd.api.types.is_float_dtype(df[col]):
                            if df[col].isna().any():
                                pass
                    
                    # Generate output file name
                    # Clean sheet name to remove invalid file system characters
                    safe_sheet_name = re.sub(r'[<>:"/\\|?*]', '_', sheet_name)
                    
                    # If only one sheet, use simple naming; otherwise include sheet name
                    if sheet_count == 1:
                        new_file = f"{base_name}.parquet"
                    else:
                        new_file = f"{base_name}_{safe_sheet_name}.parquet"
                    new_file = _non_conflicting_output_path(new_file)
                    
                    # Save to parquet
                    df.to_parquet(new_file, index=False, engine="pyarrow")
                    result_files.append(new_file)
                    
                    log_task_info(
                        task_uid, 
                        f"Sheet '{sheet_name}' converted successfully: {new_file} "
                        f"({len(df)} rows, {len(df.columns)} columns)"
                    )
                    
                except Exception as sheet_error:
                    log_task_error(task_uid, f"Failed to convert sheet '{sheet_name}': {sheet_error}")
                    # Continue processing other sheets even if one fails
                    continue
            
            # Clean up source file
            os.remove(file_path)
            
            # Return result
            if len(result_files) == 0:
                return {
                    "from": file_path,
                    "to": None,
                    "status": "failure",
                    "error": "No sheets were successfully converted",
                    "sheets_count": 0
                }
            
            log_task_info(
                task_uid, 
                f"Excel conversion completed: {len(result_files)}/{sheet_count} sheets converted successfully"
            )
            
            # Return format compatible with both single and multiple files
            return {
                "from": file_path,
                "to": result_files[0] if len(result_files) == 1 else result_files,
                "to_files": result_files,  # Always provide list for consistency
                "status": "success",
                "sheets_count": len(result_files)
            }
            
        except Exception as e:
            log_task_error(task_uid, f"convert file {file_path} error: {e}")
            return {
                "from": file_path,
                "to": None,
                "status": "failure",
                "error": str(e)
            }
    return None


def _convert_excel_to_parquet_streaming(file_path: str, task_uid, chunk_size: int) -> Optional[Dict[str, str]]:
    """
    Streaming implementation with chunked reading and writing.
    
    Note: pandas.read_excel() does NOT support chunksize parameter.
    We use openpyxl's read_only mode as a workaround.
    """
    try:
        import pyarrow as pa
        import pyarrow.parquet as pq
    except ImportError:
        log_task_error(task_uid, "PyArrow is required for streaming Parquet conversion")
        return {"from": file_path, "to": None, "status": "failure", "error": "PyArrow not installed"}
    
    if file_path.lower().endswith(".csv"):
        log_task_info(task_uid, f"[Streaming Mode] Source file address: {file_path}")
        log_task_info(task_uid, f"[Streaming Mode] Chunk size: {chunk_size:,} rows")
        
        new_file = None
        try:
            new_file = _non_conflicting_output_path(
                f"{os.path.splitext(file_path)[0]}.parquet"
            )
            
            # Count total rows
            log_task_info(task_uid, "Counting total rows...")
            total_rows = _count_csv_rows_fast(file_path)
            log_task_info(task_uid, f"Total rows: {total_rows:,}")
            
            # Stream processing
            writer = None
            rows_processed = 0
            
            for chunk in _read_csv_chunked(file_path, chunk_size):
                # Data type processing
                for col in chunk.columns:
                    if chunk[col].dtype == "object":
                        chunk[col] = chunk[col].astype(str).replace("nan", None)
                    elif pd.api.types.is_integer_dtype(chunk[col]) and chunk[col].isna().any():
                        chunk[col] = chunk[col].astype(str)
                
                # Convert to Arrow Table
                table = pa.Table.from_pandas(chunk)
                
                # Initialize or append
                if writer is None:
                    writer = pq.ParquetWriter(new_file, table.schema)
                writer.write_table(table)
                
                rows_processed += len(chunk)
                progress = (rows_processed / total_rows * 100) if total_rows > 0 else 0
                log_task_info(
                    task_uid,
                    f"Progress: {rows_processed:,}/{total_rows:,} rows ({progress:.1f}%)"
                )
                
                del chunk, table
            
            if writer:
                writer.close()
            
            return {
                "from": file_path,
                "to": new_file,
                "to_files": [new_file],
                "status": "success",
            }
        except Exception as e:
            log_task_error(task_uid, f"convert file {file_path} error: {e}")
            # Rollback
            if new_file and os.path.exists(new_file):
                try:
                    os.remove(new_file)
                except Exception:
                    pass
            return {"from": file_path, "to": None, "status": "failure", "error": str(e)}
    
    if file_path.lower().endswith((".xlsx", ".xls")):
        log_task_info(task_uid, f"[Streaming Mode] Source file address: {file_path}")
        log_task_info(task_uid, f"[Streaming Mode] Chunk size: {chunk_size:,} rows")
        
        try:
            from openpyxl import load_workbook
            
            # First pass: get sheet names
            wb_info = load_workbook(file_path, read_only=True, data_only=True)
            sheet_names = wb_info.sheetnames
            sheet_count = len(sheet_names)
            wb_info.close()
            
            log_task_info(task_uid, f"Found {sheet_count} sheet(s) in Excel file")
            
            result_files = []
            base_name = os.path.splitext(file_path)[0]
            
            for idx, sheet_name in enumerate(sheet_names, 1):
                new_file = None
                writer = None
                try:
                    log_task_info(task_uid, f"Processing sheet {idx}/{sheet_count}: '{sheet_name}'")
                    
                    safe_sheet_name = re.sub(r'[<>:"/\\|?*]', '_', sheet_name)
                    if sheet_count == 1:
                        new_file = f"{base_name}.parquet"
                    else:
                        new_file = f"{base_name}_{safe_sheet_name}.parquet"
                    new_file = _non_conflicting_output_path(new_file)
                    
                    # Count total rows using openpyxl
                    log_task_info(task_uid, f"Counting total rows in sheet '{sheet_name}'...")
                    wb_count = load_workbook(file_path, read_only=True, data_only=True)
                    ws_count = wb_count[sheet_name]
                    total_rows = ws_count.max_row - 1  # Exclude header
                    wb_count.close()
                    log_task_info(task_uid, f"Total rows: {total_rows:,}")
                    
                    # Stream processing
                    rows_processed = 0
                    
                    for chunk in _read_excel_sheet_in_chunks(file_path, sheet_name, chunk_size):
                        # Data type processing
                        for col in chunk.columns:
                            if chunk[col].dtype == "object":
                                chunk[col] = chunk[col].astype(str).replace("nan", None)
                            elif pd.api.types.is_integer_dtype(chunk[col]) and chunk[col].isna().any():
                                chunk[col] = chunk[col].astype(str)
                        
                        # Convert to Arrow Table
                        table = pa.Table.from_pandas(chunk)
                        
                        # Initialize or append
                        if writer is None:
                            writer = pq.ParquetWriter(new_file, table.schema)
                        writer.write_table(table)
                        
                        rows_processed += len(chunk)
                        progress = (rows_processed / total_rows * 100) if total_rows > 0 else 0
                        log_task_info(
                            task_uid,
                            f"Progress: {rows_processed:,}/{total_rows:,} rows ({progress:.1f}%)"
                        )
                        
                        del chunk, table
                    
                    if writer:
                        writer.close()
                    
                    result_files.append(new_file)
                    log_task_info(task_uid, f"Sheet '{sheet_name}' converted successfully to {new_file}")
                    
                except Exception as sheet_error:
                    log_task_error(task_uid, f"Failed to convert sheet '{sheet_name}': {sheet_error}")
                    # Rollback
                    if new_file and os.path.exists(new_file):
                        try:
                            os.remove(new_file)
                            log_task_info(task_uid, f"Rolled back partial file: {new_file}")
                        except Exception:
                            pass
                    continue
            
            if len(result_files) > 0:
                try:
                    os.remove(file_path)
                except PermissionError:
                    import time
                    time.sleep(0.5)
                    try:
                        os.remove(file_path)
                    except Exception as e:
                        log_task_error(task_uid, f"Warning: Could not delete source file: {e}")
            
            if len(result_files) == 0:
                return {"from": file_path, "to": None, "status": "failure", "error": "No sheets converted"}
            
            return {
                "from": file_path,
                "to": result_files[0] if len(result_files) == 1 else result_files,
                "to_files": result_files,
                "status": "success",
                "sheets_count": len(result_files)
            }
        except Exception as e:
            log_task_error(task_uid, f"convert file {file_path} error: {e}")
            return {"from": file_path, "to": None, "status": "failure", "error": str(e)}
    return None


def convert_csv_to_excel(file_path: str, task_uid, use_streaming: bool = False, chunk_size: int = 50000) -> Optional[Dict[str, str]]:
    """
    Convert CSV to Excel format.
    
    Args:
        file_path: Input file path
        task_uid: Task identifier for logging
        use_streaming: Use streaming mode for large files (default: False)
        chunk_size: Number of rows per chunk in streaming mode (default: 50000)
    
    Returns:
        Conversion result dictionary
    """
    if use_streaming:
        return _convert_csv_to_excel_streaming(file_path, task_uid, chunk_size)
    else:
        return _convert_csv_to_excel_legacy(file_path, task_uid)


def _convert_csv_to_excel_legacy(file_path: str, task_uid) -> Optional[Dict[str, str]]:
    """Original non-streaming implementation."""
    if not file_path.lower().endswith(".csv"):
        return None

    log_task_info(task_uid, f"Source file address: {file_path}")
    try:
        new_file = _non_conflicting_output_path(
            f"{os.path.splitext(file_path)[0]}.xlsx"
        )
        _read_csv(file_path).to_excel(new_file, index=False, engine="openpyxl")
        return {
            "from": file_path,
            "to": new_file,
            "to_files": [new_file],
            "status": "success",
        }
    except Exception as e:
        log_task_error(task_uid, f"convert file {file_path} error: {e}")
        return {"from": file_path, "to": None, "status": "failure", "error": str(e)}


def _convert_csv_to_excel_streaming(file_path: str, task_uid, chunk_size: int) -> Optional[Dict[str, str]]:
    """Streaming implementation using xlsxwriter's constant_memory mode."""
    if not file_path.lower().endswith(".csv"):
        return None
    
    log_task_info(task_uid, f"[Streaming Mode] Source file address: {file_path}")
    log_task_info(task_uid, f"[Streaming Mode] Chunk size: {chunk_size:,} rows")
    
    try:
        import xlsxwriter
    except ImportError:
        log_task_error(task_uid, "xlsxwriter is required for streaming Excel conversion")
        return {"from": file_path, "to": None, "status": "failure", "error": "xlsxwriter not installed"}
    
    new_file = None
    try:
        new_file = _non_conflicting_output_path(
            f"{os.path.splitext(file_path)[0]}.xlsx"
        )
        
        # Count total rows
        log_task_info(task_uid, "Counting total rows...")
        total_rows = _count_csv_rows_fast(file_path)
        log_task_info(task_uid, f"Total rows: {total_rows:,}")
        
        # Create workbook with constant_memory mode
        workbook = xlsxwriter.Workbook(new_file, {
            'constant_memory': True,
            'use_zip64': True,
            'strings_to_numbers': False,
            'strings_to_urls': False
        })
        worksheet = workbook.add_worksheet()
        
        # Stream processing
        current_row = 0
        header_written = False
        rows_processed = 0
        
        for chunk in _read_csv_chunked(file_path, chunk_size):
            # Write header (first chunk only)
            if not header_written:
                for col_idx, col_name in enumerate(chunk.columns):
                    worksheet.write(0, col_idx, col_name)
                current_row = 1
                header_written = True
            
            # Write data
            for _, data_row in chunk.iterrows():
                for col_idx, value in enumerate(data_row):
                    if pd.isna(value):
                        worksheet.write_blank(current_row, col_idx, None)
                    else:
                        worksheet.write(current_row, col_idx, value)
                current_row += 1
            
            rows_processed += len(chunk)
            progress = (rows_processed / total_rows * 100) if total_rows > 0 else 0
            log_task_info(
                task_uid,
                f"Progress: {rows_processed:,}/{total_rows:,} rows ({progress:.1f}%)"
            )
            
            del chunk
        
        # Finalize workbook
        log_task_info(task_uid, "Finalizing Excel file (building ZIP structure)...")
        workbook.close()
        
        log_task_info(task_uid, f"Conversion completed: {new_file}")
        
        return {
            "from": file_path,
            "to": new_file,
            "to_files": [new_file],
            "status": "success",
        }
    except Exception as e:
        log_task_error(task_uid, f"convert file {file_path} error: {e}")
        # Rollback
        if new_file and os.path.exists(new_file):
            try:
                os.remove(new_file)
                log_task_info(task_uid, f"Rolled back partial file: {new_file}")
            except Exception:
                pass
        return {"from": file_path, "to": None, "status": "failure", "error": str(e)}


def fix_email_links_in_html(html_content: str) -> str:
    pattern1 = r'<a href="(http://[^@]+@[^"]+?)&quot;\s*\\o\s*&quot;(http://[^"]+?)">([^<]+)</a>'

    def replace_email_link1(match):
        url1 = match.group(1)
        url2 = match.group(2)
        link_text = match.group(3)
        if url1 != url2:
            return f'<a href="{url1}">{link_text}</a> <a href="{url2}">{link_text}</a>'
        return f'<a href="{url2}">{link_text}</a>'

    pattern2 = r'<a href="(mailto:[^@]+@[^"]+?)&quot;\s*\\o\s*&quot;mailto:([^"]+?)">([^<]+)</a>'

    def replace_email_link2(match):
        address1 = match.group(1)
        address2_full = f'mailto:{match.group(2)}'
        link_text = match.group(3)
        if address1 != address2_full:
            return f'<a href="{address1}">{link_text}</a> <a href="{address2_full}">{link_text}</a>'
        return f'<a href="{address1}">{link_text}</a>'

    pattern3 = r'<a href="http://@([^"]+?)&quot;\s*\\o\s*&quot;http://([^"]+?)">([^<]+)</a>'

    def replace_email_link3(match):
        address1 = f'http://@{match.group(1)}'
        address2 = f'http://{match.group(2)}'
        link_text = match.group(3)
        if address1 != address2:
            return f'<a href="{address1}">{link_text}</a> <a href="{address2}">{link_text}</a>'
        return f'<a href="{address2}">{link_text}</a>'

    pattern4 = r'<a href="http://([^@]+@[^/"]+?)">([^<]+)</a>'

    def replace_email_link4(match):
        email = match.group(1)
        link_text = match.group(2)
        return f'<a href="http://{email}">{link_text}</a>'

    fixed_html = re.sub(pattern1, replace_email_link1, html_content)
    fixed_html = re.sub(pattern2, replace_email_link2, fixed_html)
    fixed_html = re.sub(pattern3, replace_email_link3, fixed_html)
    fixed_html = re.sub(pattern4, replace_email_link4, fixed_html)
    return fixed_html


def convert_word_to_markdown(file_path: str, task_uid) -> Optional[Dict[str, str]]:
    if file_path.lower().endswith((".docx", ".doc")):
        log_task_info(task_uid, f"Source file address：{file_path}")
        try:
            with open(file_path, "rb") as docx_file:
                result = mammoth.convert_to_html(docx_file)
                html_content = result.value
            html_content = fix_email_links_in_html(html_content)
            markdown_content = md(html_content)
            markdown_file_path = os.path.splitext(file_path)[0] + ".md"
            with open(markdown_file_path, "w", encoding="utf-8") as md_file:
                md_file.write(markdown_content)
            log_task_info(task_uid, f"convert file {markdown_file_path} succeed")
            os.remove(file_path)
            return {"from": file_path, "to": markdown_file_path, "status": "success"}
        except Exception as e:
            log_task_error(task_uid, f"convert file {file_path} error: {e}")
            return {"from": file_path, "to": None, "status": "failure"}
    return None


def _read_text_file_bytes(raw: bytes) -> str:
    for enc in ("utf-8-sig", "utf-8", "gbk", "gb18030", "latin-1"):
        try:
            return raw.decode(enc)
        except UnicodeDecodeError:
            continue
    return raw.decode("utf-8", errors="replace")


def convert_txt_to_markdown(file_path: str, task_uid) -> Optional[Dict[str, str]]:
    if not file_path.lower().endswith(".txt"):
        return None
    log_task_info(task_uid, f"Source file address：{file_path}")
    try:
        with open(file_path, "rb") as f:
            raw = f.read()
        text_content = _read_text_file_bytes(raw)
        markdown_file_path = os.path.splitext(file_path)[0] + ".md"
        with open(markdown_file_path, "w", encoding="utf-8") as md_file:
            md_file.write(text_content)
        log_task_info(task_uid, f"convert file {markdown_file_path} succeed")
        os.remove(file_path)
        return {"from": file_path, "to": markdown_file_path, "status": "success"}
    except Exception as e:
        log_task_error(task_uid, f"convert file {file_path} error: {e}")
        return {"from": file_path, "to": None, "status": "failure"}


def convert_html_to_markdown(file_path: str, task_uid) -> Optional[Dict[str, str]]:
    if not file_path.lower().endswith((".html", ".htm")):
        return None
    log_task_info(task_uid, f"Source file address：{file_path}")
    try:
        with open(file_path, "rb") as f:
            raw = f.read()
        html_content = _read_text_file_bytes(raw)
        html_content = fix_email_links_in_html(html_content)
        markdown_content = md(html_content)
        markdown_file_path = os.path.splitext(file_path)[0] + ".md"
        with open(markdown_file_path, "w", encoding="utf-8") as md_file:
            md_file.write(markdown_content)
        log_task_info(task_uid, f"convert file {markdown_file_path} succeed")
        os.remove(file_path)
        return {"from": file_path, "to": markdown_file_path, "status": "success"}
    except Exception as e:
        log_task_error(task_uid, f"convert file {file_path} error: {e}")
        return {"from": file_path, "to": None, "status": "failure"}


def convert_ppt_to_markdown(file_path: str, task_uid) -> Optional[Dict[str, str]]:
    if file_path.lower().endswith((".pptx", ".ppt")):
        log_task_info(task_uid, f"Source file address：{file_path}")
        try:
            prs = Presentation(file_path)
            markdown_content = ""
            for i, slide in enumerate(prs.slides):
                markdown_content += f" lantern slide {i + 1}\n\n"
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text_content = shape.text.strip()
                        if text_content:
                            if len(text_content) < 50 and "\n" not in text_content:
                                markdown_content += f"## {text_content}\n\n"
                            else:
                                markdown_content += f"{text_content}\n\n"
            markdown_file_path = os.path.splitext(file_path)[0] + ".md"
            with open(markdown_file_path, "w", encoding="utf-8") as md_file:
                md_file.write(markdown_content)
            log_task_info(task_uid, f"convert file {markdown_file_path} succeed")
            os.remove(file_path)
            return {"from": file_path, "to": markdown_file_path, "status": "success"}
        except Exception as e:
            log_task_error(task_uid, f"convert file {file_path} error: {e}")
            return {"from": file_path, "to": None, "status": "failure"}
    return None


def convert_pdf_to_markdown(
    file_path: str,
    task_uid,
    mineru_api_url: Optional[str] = None,
    mineru_backend: Optional[str] = None,
) -> Optional[Dict[str, str]]:
    if file_path.lower().endswith(".pdf"):
        log_task_info(task_uid, f"Source file address：{file_path}")
        try:
            from mineru.cli.common import prepare_env
            from mineru.data.data_reader_writer import FileBasedDataWriter
            from mineru.backend.vlm.vlm_middle_json_mkcontent import union_make as vlm_union_make
            from mineru.utils.enum_class import MakeMode

            server_url = mineru_api_url or os.getenv("MINERU_API_URL", "http://111.4.242.20:30000")
            backend = mineru_backend or os.getenv("MINERU_BACKEND", "http-client")
            log_task_info(task_uid, f"Using MinerU API server: {server_url}, backend: {backend}")

            pdf_file_name = Path(file_path).stem
            temp_output_dir = Path(file_path).parent / f"_temp_pdf_convert_{pdf_file_name}"
            temp_output_dir.mkdir(exist_ok=True)

            result_json_path = temp_output_dir / "mineru_result.json"
            repo_root = Path(__file__).resolve().parents[2]
            mineru_worker_script = repo_root / "data_server" / "pod" / "mineru_worker.py"

            cmd = [
                sys.executable,
                str(mineru_worker_script),
                file_path,
                str(temp_output_dir),
                server_url,
                backend,
                str(result_json_path),
            ]
            process = subprocess.Popen(cmd, cwd=str(repo_root))
            process.wait()

            if process.returncode != 0:
                error_msg = "MinerU subprocess failed"
                if result_json_path.exists():
                    try:
                        with open(result_json_path, "r", encoding="utf-8") as f:
                            result_data = json.load(f)
                            if not result_data.get("success", False):
                                error_msg = result_data.get("error", error_msg)
                    except Exception:
                        pass
                log_task_error(task_uid, f"MinerU subprocess failed: {error_msg}")
                raise RuntimeError(f"MinerU subprocess failed: {error_msg}")

            if not result_json_path.exists():
                raise FileNotFoundError(f"Result JSON file not found: {result_json_path}")

            with open(result_json_path, "r", encoding="utf-8") as f:
                result_data = json.load(f)

            if not result_data.get("success", False):
                error_msg = result_data.get("error", "Unknown error")
                log_task_error(task_uid, f"MinerU subprocess error: {error_msg}")
                raise RuntimeError(f"MinerU subprocess error: {error_msg}")

            middle_json = result_data["middle_json"]
            local_image_dir, local_md_dir = prepare_env(str(temp_output_dir), pdf_file_name, "vlm")
            _ = local_image_dir
            md_writer = FileBasedDataWriter(local_md_dir)

            pdf_info = middle_json["pdf_info"]
            image_dir = ""
            md_content_str = vlm_union_make(pdf_info, MakeMode.MM_MD, image_dir)

            markdown_filename = f"{pdf_file_name}.md"
            md_writer.write_string(markdown_filename, md_content_str)

            markdown_file_path = Path(local_md_dir) / markdown_filename
            final_markdown_path = os.path.splitext(file_path)[0] + ".md"
            shutil.move(str(markdown_file_path), final_markdown_path)

            log_task_info(task_uid, f"convert file {final_markdown_path} succeed")
            os.remove(file_path)
            if temp_output_dir.exists():
                shutil.rmtree(temp_output_dir)
            return {"from": file_path, "to": final_markdown_path, "status": "success"}
        except Exception as e:
            log_task_error(task_uid, f"convert file {file_path} error: {e}")
            temp_output_dir = Path(file_path).parent / f"_temp_pdf_convert_{Path(file_path).stem}"
            if temp_output_dir.exists():
                try:
                    shutil.rmtree(temp_output_dir)
                except Exception:
                    pass
            return {"from": file_path, "to": None, "status": "failure"}
    return None


def search_files(folder_path: str, types: List[int]) -> Tuple[bool, List[str]]:
    type_map: Dict[int, List[str]] = {
        0: [".ppt", ".pptx"],
        1: [".doc", ".docx"],
        # Excel tasks also accept CSV files in mixed datasets.
        3: [".xls", ".xlsx", ".csv"],
        5: [".csv"],
        7: [".pdf"],
        8: [".txt"],
        9: [".html", ".htm"],
    }

    target_extensions = set()
    for file_type in types:
        if file_type in type_map:
            for ext in type_map[file_type]:
                target_extensions.add(ext.lower())

    found_files: List[str] = []

    def traverse(current_path: str) -> None:
        try:
            entries = os.listdir(current_path)
            for entry in entries:
                entry_path = os.path.join(current_path, entry)
                if os.path.isdir(entry_path):
                    traverse(entry_path)
                elif os.path.isfile(entry_path):
                    file_ext = os.path.splitext(entry)[1].lower()
                    if file_ext in target_extensions:
                        found_files.append(entry_path)
        except PermissionError:
            print(f"No permission to access the folder: {current_path}")
        except Exception as e:
            print(f"Processing path {current_path} error: {str(e)}")

    traverse(folder_path)
    return bool(len(found_files) > 0), found_files
