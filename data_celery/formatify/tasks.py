import os
import shutil
import time
import traceback
import json
import subprocess
import sys
import re
from typing import List, Dict, Tuple, Optional
from data_celery.main import celery_app
from sqlalchemy.orm import Session
from loguru import logger
from pathlib import Path

from data_celery.utils import get_format_folder_path, ensure_directory_exists, ensure_directory_exists_remove
from data_celery.db import FormatifyManager
from data_engine.exporter.load import load_exporter
from data_server.formatify.FormatifyModels import DataFormatTask, DataFormatTaskStatusEnum, DataFormatTypeEnum, \
    getFormatTypeName
from data_server.database.session import get_sync_session
from sqlalchemy import text
from data_engine.ingester.load import load_ingester
import pandas as pd
import mammoth
from markdownify import markdownify as md
from pptx import Presentation
from data_celery.mongo_tools.tools import insert_formatity_task_log_info, insert_formatity_task_log_error
from pycsghub.upload_large_folder.main import upload_large_folder_internal
from pycsghub.cmd.repo_types import RepoType
from pycsghub.utils import (build_csg_headers,
                            model_id_to_group_owner_name,
                            get_endpoint,
                            REPO_TYPE_DATASET)
from data_engine.utils.env import GetHubEndpoint

def _refresh_task_with_reconnect(db_session: Session, format_task: DataFormatTask, task_id: int) -> Tuple[Session, DataFormatTask]:
    """
    Refresh task status, automatically reconnect if database connection fails
    
    Returns:
        Tuple[Session, DataFormatTask]: New session and task object
    """
    try:
        db_session.refresh(format_task)
        return db_session, format_task
    except Exception as db_error:
        logger.warning(f"Database connection error during refresh: {db_error}, attempting to reconnect")
        try:
            db_session.rollback()
        except:
            pass
        try:
            db_session.close()
        except:
            pass
        # Re-acquire session and task object
        new_session = get_sync_session()
        new_task = FormatifyManager.get_formatify_task(new_session, task_id)
        return new_session, new_task

def _commit_with_retry(db_session: Session, format_task: DataFormatTask, task_id: int, task_status: int):
    """
    Commit database changes, retry with new session if failed
    
    Args:
        db_session: Current database session
        format_task: Task object
        task_id: Task ID
        task_status: Task status to set
    """
    try:
        format_task.task_status = task_status
        db_session.commit()
    except Exception as commit_error:
        logger.error(f"Failed to commit task status: {commit_error}")
        try:
            db_session.rollback()
            db_session.close()
        except:
            pass
        # Try to commit with new session
        try:
            new_session = get_sync_session()
            retry_task = FormatifyManager.get_formatify_task(new_session, task_id)
            retry_task.task_status = task_status
            new_session.commit()
            new_session.close()
        except Exception as retry_error:
            logger.error(f"Failed to retry commit: {retry_error}")

@celery_app.task
def format_task(task_id: int, user_name: str, user_token: str):

    tmp_path: str = None
    db_session: Session = None
    format_task: DataFormatTask = None
    task_uid: str = None  # Save task_uid to avoid accessing database object in finally block

    try:
        db_session: Session = get_sync_session()
        format_task: DataFormatTask = FormatifyManager.get_formatify_task(db_session, task_id)
        task_uid = format_task.task_uid  # Save to local variable
        tmp_path = get_format_folder_path(task_uid)
        insert_formatity_task_log_info(task_uid, f"Create temporary directory：{tmp_path}")
        ensure_directory_exists(tmp_path)

        insert_formatity_task_log_info(task_uid, f"Start downloading directory....")
        try:
            ingesterCSGHUB = load_ingester(
                dataset_path=tmp_path,
                repo_id=format_task.from_csg_hub_repo_id,
                branch=format_task.from_csg_hub_dataset_branch,
                user_name=user_name,
                user_token=user_token,
            )
            ingester_result = ingesterCSGHUB.ingest()
            insert_formatity_task_log_info(task_uid, f"Download directory completed... Directory address：{ingester_result}")
        except Exception as e:
            error_msg = str(e)
            # Check if it's an authentication error
            if "401" in error_msg or "Unauthorized" in error_msg:
                detailed_error = f"认证失败：无法访问数据集 {format_task.from_csg_hub_repo_id}。可能原因：1) Token 已过期，请重新登录；2) Token 无效；3) 没有访问该数据集的权限。错误详情：{error_msg}"
                insert_formatity_task_log_error(task_uid, detailed_error)
                logger.error(f"Authentication failed for task {task_uid}: {error_msg}")
            else:
                insert_formatity_task_log_error(task_uid, f"Failed to download dataset: {error_msg}")
                logger.error(f"Failed to download dataset for task {task_uid}: {error_msg}")
            _commit_with_retry(db_session, format_task, task_id, DataFormatTaskStatusEnum.ERROR.value)
            return
        work_dir = Path(tmp_path).joinpath('work')
        file_bool = search_files(tmp_path,[format_task.from_data_type])

        if not file_bool:
            insert_formatity_task_log_info(task_uid, f"file not found. task ended....")
            _commit_with_retry(db_session, format_task, task_id, DataFormatTaskStatusEnum.ERROR.value)
            return
        insert_formatity_task_log_info(task_uid, f"Start converting file...")
        
        # Create new directory containing only successfully converted files
        converted_dir = Path(tmp_path).joinpath('converted')
        ensure_directory_exists(str(converted_dir))
        
        # Check if skip_meta is enabled
        # Note: skip_meta=True means upload meta file, skip_meta=False means skip uploading meta file
        # Get skip_meta from database
        skip_meta = False  # Default value
        try:
            # Refresh the object to ensure we have latest data from database
            db_session.refresh(format_task)
            # Get skip_meta value directly from the object
            skip_meta_value = getattr(format_task, 'skip_meta', None)
            if skip_meta_value is not None:
                skip_meta = bool(skip_meta_value)
            else:
                # If None, try raw SQL query as fallback
                result = db_session.execute(text("""
                    SELECT skip_meta 
                    FROM data_format_tasks 
                    WHERE id = :task_id
                """), {"task_id": format_task.id})
                row = result.fetchone()
                if row and row[0] is not None:
                    skip_meta = bool(row[0])
                else:
                    skip_meta = False
        except Exception as e:
            # Column doesn't exist in database or query failed, default to False (don't upload meta file)
            logger.warning(f"Could not read skip_meta from database: {e}, defaulting to False")
            skip_meta = False
        
        insert_formatity_task_log_info(task_uid, f'skip_meta value: {skip_meta} (True=upload meta.log, False=skip meta.log)')
        
        # Create meta folder (only if skip_meta is True, meaning we should upload meta file)
        meta_dir = None
        if skip_meta:
            meta_dir = converted_dir / "meta"
            ensure_directory_exists(str(meta_dir))
            insert_formatity_task_log_info(task_uid, f'Creating meta.log file and meta folder (skip_meta=True)')
            insert_formatity_task_log_info(task_uid, f'Created meta directory: {meta_dir}')
        else:
            insert_formatity_task_log_info(task_uid, 'Skipping meta.log file and meta folder creation (skip_meta=False)')
        
        # Get target file extension list
        type_map: Dict[int, List[str]] = {
            0: ['.ppt', '.pptx'],  # PPT
            1: ['.doc', '.docx'],  # Word
            3: ['.xls', '.xlsx'],  # Excel
            7: ['.pdf']  # PDF
        }
        target_extensions = set()
        if format_task.from_data_type in type_map:
            for ext in type_map[format_task.from_data_type]:
                target_extensions.add(ext.lower())
        
        # First count total files (using exactly the same logic as conversion)
        total_count = 0
        base_path = Path(ingester_result)
        file_list_for_count = []
        for root, dirs, files in os.walk(ingester_result):
            for file in files:
                file_path_full = os.path.join(root, file)
                # Only process files of specified types (exactly the same logic as conversion)
                file_ext = Path(file_path_full).suffix.lower()
                if file_ext in target_extensions:
                    total_count += 1
                    file_list_for_count.append(file_path_full)
        
        insert_formatity_task_log_info(task_uid, f'Found {total_count} files to convert')
        if total_count > 0:
            insert_formatity_task_log_info(task_uid, f'Files to convert: {[Path(f).name for f in file_list_for_count[:5]]}{"..." if total_count > 5 else ""}')
        
        # Initialize statistics
        success_count = 0
        failure_count = 0
        
        # Initialize meta.log data structure (only if skip_meta is True, meaning we should upload meta file)
        meta_data = None
        meta_file_path = None
        if skip_meta:
            # Note: total is set during initialization and not modified afterwards, only success and failure are updated
            meta_data = {
                "job_id": format_task.id,
                "job_name": format_task.name,
                "source_repo": format_task.from_csg_hub_repo_id,
                "source_branch": format_task.from_csg_hub_dataset_branch,
                "files": [],
                "result": {
                    "total": total_count,  # Fixed total, not modified afterwards
                    "success": 0,
                    "failure": 0
                }
            }
            meta_file_path = meta_dir / "meta.log"
            
            # First upload meta.log file containing total count
            with open(meta_file_path, 'w', encoding='utf-8') as f:
                json.dump(meta_data, f, indent=2, ensure_ascii=False)
            insert_formatity_task_log_info(task_uid, f'Generated initial meta.log file with total: {total_count} (total will remain fixed)')
        
        # Create exporter (reuse the same instance)
        exporter = load_exporter(
            export_path=str(converted_dir),
            repo_id=format_task.to_csg_hub_repo_id,
            branch=format_task.to_csg_hub_dataset_default_branch,
            user_name=user_name,
            user_token=user_token,
            path_is_dir=True,
            work_dir=str(work_dir)
        )
        
        # First upload meta.log containing total count (only if skip_meta is True, meaning we should upload meta file)
        # If it fails 3 times, the task terminates immediately
        # Note: First attempt doesn't count in retry count, retry counting only starts after failure
        initial_upload_success = True  # Default to True if skipping meta (skip_meta=False)
        if skip_meta:
            initial_upload_success = False
            initial_upload_retry_count = 0
            max_initial_retries = 3  # Maximum 3 retries for initial upload
            
            # First attempt (doesn't count in retry count)
            try:
                insert_formatity_task_log_info(task_uid, f'开始上传初始 meta.log (总文件数: {total_count})')
                exporter.export_large_folder()
                insert_formatity_task_log_info(task_uid, f'[成功] 初始 meta.log 上传成功 (总文件数: {total_count})')
                initial_upload_success = True
                upload_failure_count = 0  # Reset failure count
            except Exception as e:
                # First failure, start retrying
                initial_upload_retry_count += 1
                error_msg = f'[上传失败 {initial_upload_retry_count}/{max_initial_retries}] 初始 meta.log 上传失败: {str(e)}'
                insert_formatity_task_log_error(task_uid, error_msg)
                logger.error(error_msg)
                
                # Retry up to 3 times
                while not initial_upload_success and initial_upload_retry_count < max_initial_retries:
                    try:
                        initial_upload_retry_count += 1
                        insert_formatity_task_log_info(task_uid, f'[重试 {initial_upload_retry_count}/{max_initial_retries}] 开始上传初始 meta.log (总文件数: {total_count})')
                        exporter.export_large_folder()
                        insert_formatity_task_log_info(task_uid, f'[成功] 初始 meta.log 上传成功 (总文件数: {total_count})')
                        initial_upload_success = True
                        upload_failure_count = 0  # Reset failure count
                    except Exception as e:
                        error_msg = f'[上传失败 {initial_upload_retry_count}/{max_initial_retries}] 初始 meta.log 上传失败: {str(e)}'
                        insert_formatity_task_log_error(task_uid, error_msg)
                        logger.error(error_msg)
                        
                        if initial_upload_retry_count >= max_initial_retries:
                            # Initial upload failed 3 times, task terminates immediately
                            final_error_msg = f'[任务终止] 初始 meta.log 上传失败 {max_initial_retries} 次，任务已停止。总文件数: {total_count}'
                        insert_formatity_task_log_error(task_uid, final_error_msg)
                        logger.error(final_error_msg)
                        _commit_with_retry(db_session, format_task, task_id, DataFormatTaskStatusEnum.ERROR.value)
                        raise RuntimeError(f'初始 meta.log 上传失败 {max_initial_retries} 次，任务已终止。总文件数: {total_count}')
        
        # If initial upload fails, should not continue execution
        if not initial_upload_success:
            error_msg = f'[任务终止] 初始 meta.log 上传失败，任务已停止。总文件数: {total_count}'
            insert_formatity_task_log_error(task_uid, error_msg)
            logger.error(error_msg)
            _commit_with_retry(db_session, format_task, task_id, DataFormatTaskStatusEnum.ERROR.value)
            raise RuntimeError(f'初始 meta.log 上传失败，任务已终止。总文件数: {total_count}')
        
        # Select conversion function
        convert_func = None
        match format_task.from_data_type:
            case DataFormatTypeEnum.Excel.value:
                match format_task.to_data_type:
                    case DataFormatTypeEnum.Csv.value:
                        convert_func = convert_excel_to_csv
                    case DataFormatTypeEnum.Json.value:
                        convert_func = convert_excel_to_json
                    case DataFormatTypeEnum.Parquet.value:
                        convert_func = convert_excel_to_parquet
            case DataFormatTypeEnum.Word.value:
                match format_task.to_data_type:
                    case DataFormatTypeEnum.Markdown.value:
                        convert_func = convert_word_to_markdown
            case DataFormatTypeEnum.PPT.value:
                match format_task.to_data_type:
                    case DataFormatTypeEnum.Markdown.value:
                        convert_func = convert_ppt_to_markdown
            case DataFormatTypeEnum.PDF.value:
                match format_task.to_data_type:
                    case DataFormatTypeEnum.Markdown.value:
                        convert_func = convert_pdf_to_markdown
        
        if convert_func is None:
            insert_formatity_task_log_error(task_uid, f"Unsupported conversion: {getFormatTypeName(format_task.from_data_type)} -> {getFormatTypeName(format_task.to_data_type)}")
            _commit_with_retry(db_session, format_task, task_id, DataFormatTaskStatusEnum.ERROR.value)
            return
        
        insert_formatity_task_log_info(task_uid,
                                       f"Change the table of contents：{ingester_result}，Source file type：{getFormatTypeName(format_task.from_data_type)}，Target file type：{getFormatTypeName(format_task.to_data_type)}")
        
        # Track used filenames to handle duplicate names
        used_names = {}
        
        # Iterate files and convert/upload in real-time
        upload_failure_count = 0  # Consecutive upload failure count (initial upload failure already counted)
        max_upload_failures = 3  # Maximum consecutive upload failures
        
        for root, dirs, files in os.walk(ingester_result):
            # Check if task is stopped - with database connection error handling
            db_session, format_task = _refresh_task_with_reconnect(db_session, format_task, task_id)
            if format_task.task_status == DataFormatTaskStatusEnum.STOP.value:
                insert_formatity_task_log_info(task_uid, "Task stopped by user")
                return
            
            for file in files:
                # Check if task is stopped - with database connection error handling
                db_session, format_task = _refresh_task_with_reconnect(db_session, format_task, task_id)
                if format_task.task_status == DataFormatTaskStatusEnum.STOP.value:
                    insert_formatity_task_log_info(task_uid, "Task stopped by user")
                    return
                
                file_path_full = os.path.join(root, file)
                
                # Only process files of specified types
                file_ext = Path(file_path_full).suffix.lower()
                if file_ext not in target_extensions:
                    continue
                
                # Execute conversion
                # For PDF to MD conversion, need to pass mineru_api_url and mineru_backend parameters
                if format_task.from_data_type == DataFormatTypeEnum.PDF.value and format_task.to_data_type == DataFormatTypeEnum.Markdown.value:
                    result = convert_func(file_path_full, task_uid, format_task.mineru_api_url, format_task.mineru_backend)
                else:
                    result = convert_func(file_path_full, task_uid)
                
                if not isinstance(result, dict):
                    continue
                
                from_file = result['from']
                to_file = result['to']
                status = result['status']
                
                # Calculate relative path of original file (relative to base_path) for meta.log
                try:
                    from_rel_path = str(Path(from_file).relative_to(base_path))
                except ValueError:
                    from_rel_path = Path(from_file).name
                
                # Ensure from path uses forward slashes (cross-platform compatibility)
                from_rel_path = from_rel_path.replace('\\', '/')
                
                if status == 'success' and to_file and Path(to_file).exists():
                    src_path = Path(to_file)
                    # Use filename directly, add sequence number if duplicate
                    file_name = src_path.name
                    if file_name in used_names:
                        # Filename conflict, add sequence number
                        name_part = src_path.stem
                        ext_part = src_path.suffix
                        counter = used_names[file_name]
                        new_name = f"{name_part}_{counter}{ext_part}"
                        used_names[file_name] += 1
                        used_names[new_name] = 1
                        dst_path = converted_dir / new_name
                        final_to_file = new_name
                        insert_formatity_task_log_info(task_uid, f'Copied converted file (renamed due to conflict): {file_name} -> {new_name}')
                    else:
                        dst_path = converted_dir / file_name
                        used_names[file_name] = 1
                        final_to_file = file_name
                        insert_formatity_task_log_info(task_uid, f'Copied converted file: {file_name}')
                    
                    # Copy file to converted directory
                    shutil.copy2(src_path, dst_path)
                    
                    # to path: relative to converted_dir
                    to_rel_path = final_to_file.replace('\\', '/')
                    
                    # Immediately upload successfully converted file (update statistics only after successful upload)
                    try:
                        exporter.export_large_folder()
                        insert_formatity_task_log_info(task_uid, f'Uploaded converted file: {final_to_file}')
                        
                        # Update statistics only after file upload succeeds
                        success_count += 1
                        if skip_meta:
                            meta_data["files"].append({
                                "from": from_rel_path,
                                "to": to_rel_path,
                                "status": "success"
                            })
                            meta_data["result"]["success"] = success_count
                        upload_failure_count = 0  # Reset failure count
                    except Exception as e:
                        upload_failure_count += 1
                        error_msg = f'[上传失败 {upload_failure_count}/{max_upload_failures}] 文件上传失败: {final_to_file}, 错误: {str(e)}'
                        insert_formatity_task_log_error(task_uid, error_msg)
                        logger.error(error_msg)
                        
                        # File upload failed, immediately update statistics to failure and update meta.log
                        failure_count += 1
                        if skip_meta:
                            meta_data["files"].append({
                                "from": from_rel_path,
                                "to": to_rel_path,
                                "status": "failure",
                                "error": f"上传失败: {str(e)}"
                            })
                            meta_data["result"]["failure"] = failure_count
                            
                            # Ensure total remains unchanged (should not be modified)
                            assert meta_data["result"]["total"] == total_count, f"Total should remain {total_count}, but got {meta_data['result']['total']}"
                            
                            # Immediately update and upload meta.log (record failed uploads)
                            with open(meta_file_path, 'w', encoding='utf-8') as f:
                                json.dump(meta_data, f, indent=2, ensure_ascii=False)
                            try:
                                exporter.export_large_folder()
                                insert_formatity_task_log_info(task_uid, f'Updated meta.log with upload failure record (total: {total_count}, success: {success_count}, failure: {failure_count})')
                                # Note: meta.log upload success doesn't reset upload_failure_count, only file upload success resets it
                            except Exception as e:
                                # meta.log update/upload failure only logs, doesn't raise exception, doesn't affect file upload failure count
                                error_msg = f'meta.log 更新上传失败 (文件: {final_to_file}), 错误: {str(e)}'
                                insert_formatity_task_log_error(task_uid, error_msg)
                                logger.error(error_msg)
                                # Don't increment upload_failure_count, doesn't affect file upload failure count
                        
                        # If too many consecutive upload failures, stop task
                        if upload_failure_count >= max_upload_failures:
                            error_msg = f'[任务终止] 连续上传失败 {upload_failure_count} 次，任务已停止。'
                            insert_formatity_task_log_error(task_uid, error_msg)
                            logger.error(error_msg)
                            _commit_with_retry(db_session, format_task, task_id, DataFormatTaskStatusEnum.ERROR.value)
                            raise RuntimeError(f'连续上传失败 {upload_failure_count} 次，任务已停止')
                    
                    # When file upload succeeds, update and upload meta.log (only if skip_meta is True)
                    if skip_meta:
                        # Ensure total remains unchanged (should not be modified)
                        assert meta_data["result"]["total"] == total_count, f"Total should remain {total_count}, but got {meta_data['result']['total']}"
                        
                        if upload_failure_count == 0:  # Only update when upload succeeds (already updated above on failure)
                            with open(meta_file_path, 'w', encoding='utf-8') as f:
                                json.dump(meta_data, f, indent=2, ensure_ascii=False)
                            try:
                                exporter.export_large_folder()
                                insert_formatity_task_log_info(task_uid, f'Updated and uploaded meta.log (total: {total_count}, success: {success_count}, failure: {failure_count})')
                                # Note: meta.log upload success doesn't reset upload_failure_count, only file upload success resets it
                            except Exception as e:
                                # meta.log update/upload failure only logs, doesn't raise exception, doesn't affect file upload failure count
                                error_msg = f'meta.log 更新上传失败 (文件: {final_to_file}), 错误: {str(e)}'
                                insert_formatity_task_log_error(task_uid, error_msg)
                                logger.error(error_msg)
                                # Don't increment upload_failure_count, doesn't affect file upload failure count
                else:
                    # Also record failed conversions to meta.log (only if skip_meta is True)
                    failure_count += 1
                    if skip_meta:
                        meta_entry = {
                            "from": from_rel_path,
                            "to": None,
                            "status": "failure"
                        }
                        # If there's error message, also record to meta.log
                        if "error" in result:
                            meta_entry["error"] = result["error"]
                        meta_data["files"].append(meta_entry)
                        meta_data["result"]["failure"] = failure_count
                        # Ensure total remains unchanged (should not be modified)
                        assert meta_data["result"]["total"] == total_count, f"Total should remain {total_count}, but got {meta_data['result']['total']}"
                        
                        # Update and upload meta.log (include failure records, don't raise exception, only log)
                        with open(meta_file_path, 'w', encoding='utf-8') as f:
                            json.dump(meta_data, f, indent=2, ensure_ascii=False)
                        try:
                            exporter.export_large_folder()
                            insert_formatity_task_log_info(task_uid, f'Updated meta.log with failure record (total: {total_count}, success: {success_count}, failure: {failure_count})')
                            # Note: meta.log upload success doesn't reset upload_failure_count, only file upload success resets it
                        except Exception as e:
                            # meta.log update/upload failure only logs, doesn't raise exception, doesn't affect file upload failure count
                            error_msg = f'meta.log 更新上传失败 (转换失败的文件), 错误: {str(e)}'
                            insert_formatity_task_log_error(task_uid, error_msg)
                            logger.error(error_msg)
                            # Don't increment upload_failure_count, doesn't affect file upload failure count
        
        insert_formatity_task_log_info(task_uid, f'All files processed. Total: {total_count}, Success: {success_count}, Failure: {failure_count}')
        _commit_with_retry(db_session, format_task, task_id, DataFormatTaskStatusEnum.COMPLETED.value)
        pass
    except Exception as e:
        traceback.print_exc()
        # Use saved task_uid to avoid accessing potentially invalid database object
        if task_uid:
            insert_formatity_task_log_error(task_uid, f"Conversion task failed: {str(e)}")
        
        # Try to update task status to ERROR
        try:
            if db_session and format_task:
                _commit_with_retry(db_session, format_task, task_id, DataFormatTaskStatusEnum.ERROR.value)
        except Exception as db_error:
            logger.error(f"Failed to update task status to ERROR: {db_error}")
            # If current session fails, try with new session
            try:
                new_session = get_sync_session()
                retry_task = FormatifyManager.get_formatify_task(new_session, task_id)
                retry_task.task_status = DataFormatTaskStatusEnum.ERROR.value
                new_session.commit()
                new_session.close()
            except Exception as retry_error:
                logger.error(f"Failed to retry update task status: {retry_error}")
    finally:
        # Safely cleanup resources
        if db_session:
            try:
                db_session.close()
            except:
                pass

        if tmp_path:
            try:
                shutil.rmtree(tmp_path)
                if task_uid:  # Use saved task_uid to avoid accessing database object
                    insert_formatity_task_log_info(task_uid, f"Delete temporary directory：{tmp_path}")
            except Exception as cleanup_error:
                logger.error(f"Failed to cleanup temporary directory {tmp_path}: {cleanup_error}")


def format_task_func(
        tmp_path: str,
        from_type: int,
        to_type: int,
        task_uid: str,
        mineru_api_url: Optional[str] = None
) -> List[Dict[str, str]]:
    """
    Execute format conversion task
    Returns a list of conversion results, each containing from, to, status
    """
    insert_formatity_task_log_info(task_uid,
                                   f"Change the table of contents：{tmp_path}，Source file type：{getFormatTypeName(from_type)}，Target file type：{getFormatTypeName(to_type)}")
    conversion_results = []
    match from_type:
        case DataFormatTypeEnum.Excel.value:
            match to_type:
                case DataFormatTypeEnum.Csv.value:
                    conversion_results = traverse_files(tmp_path, convert_excel_to_csv, task_uid)
                case DataFormatTypeEnum.Json.value:
                    conversion_results = traverse_files(tmp_path, convert_excel_to_json, task_uid)
                case DataFormatTypeEnum.Parquet.value:
                    conversion_results = traverse_files(tmp_path, convert_excel_to_parquet, task_uid)
        case DataFormatTypeEnum.Word.value:
            match to_type:
                case DataFormatTypeEnum.Markdown.value:
                    conversion_results = traverse_files(tmp_path, convert_word_to_markdown, task_uid)
        case DataFormatTypeEnum.PPT.value:
            match to_type:
                case DataFormatTypeEnum.Markdown.value:
                    conversion_results = traverse_files(tmp_path, convert_ppt_to_markdown, task_uid)
        case DataFormatTypeEnum.PDF.value:
            match to_type:
                case DataFormatTypeEnum.Markdown.value:
                    # For PDF to MD conversion, need to pass mineru_api_url
                    conversion_results = traverse_files(tmp_path, convert_pdf_to_markdown, task_uid, mineru_api_url)
    return conversion_results


def traverse_files(file_path: str, func, task_uid, mineru_api_url: Optional[str] = None) -> List[Dict[str, str]]:
    """
    Traverse files and execute conversion function
    Returns a list of conversion results, each containing from, to, status
    """
    conversion_results = []
    for root, dirs, files in os.walk(file_path):
        for file in files:
            file_path_full = os.path.join(root, file)
            # For PDF to MD conversion, pass mineru_api_url parameter
            if func == convert_pdf_to_markdown and mineru_api_url is not None:
                result = func(file_path_full, task_uid, mineru_api_url)
            else:
                result = func(file_path_full, task_uid)
            # result should be a dictionary containing from, to, status
            if isinstance(result, dict):
                conversion_results.append(result)
        for dir in dirs:
            dir_path = os.path.join(root, dir)
            sub_results = traverse_files(dir_path, func, task_uid, mineru_api_url)
            conversion_results.extend(sub_results)
    return conversion_results


def convert_excel_to_csv(file_path: str, task_uid) -> Optional[Dict[str, str]]:
    if file_path.lower().endswith(('.xlsx', '.xls')):
        insert_formatity_task_log_info(task_uid, f'Source file address：{file_path}')
        try:
            df = pd.read_excel(file_path)
            new_file = os.path.splitext(file_path)[0] + '.csv'
            df.to_csv(new_file, index=False)
            insert_formatity_task_log_info(task_uid, f'convert file {new_file} succeed')
            os.remove(file_path)
            return {
                "from": file_path,
                "to": new_file,
                "status": "success"
            }
        except Exception as e:
            print(f"convert file {file_path} error: {e}")
            insert_formatity_task_log_error(task_uid, f"convert file {file_path} error: {e}")
            return {
                "from": file_path,
                "to": None,
                "status": "failure"
            }
    else:
        return None  # Not a target file, return None


def convert_excel_to_json(file_path: str, task_uid) -> Optional[Dict[str, str]]:
    if file_path.lower().endswith(('.xlsx', '.xls')):
        insert_formatity_task_log_info(task_uid, f'Source file address：{file_path}')
        try:
            df = pd.read_excel(file_path)
            new_file = os.path.splitext(file_path)[0] + '.json'
            df.to_json(new_file, orient='records', force_ascii=False)
            insert_formatity_task_log_info(task_uid, f'convert file {new_file} succeed')
            os.remove(file_path)
            return {
                "from": file_path,
                "to": new_file,
                "status": "success"
            }
        except Exception as e:
            print(f"convert file {file_path} error: {e}")
            insert_formatity_task_log_error(task_uid, f"convert file {file_path} error: {e}")
            return {
                "from": file_path,
                "to": None,
                "status": "failure"
            }
    else:
        return None  # Not a target file, return None


def convert_excel_to_parquet(file_path: str, task_uid) -> Optional[Dict[str, str]]:
    if file_path.lower().endswith(('.xlsx', '.xls')):
        insert_formatity_task_log_info(task_uid, f'Source file address：{file_path}')
        try:
            # Read Excel file
            df = pd.read_excel(file_path)
            
            # Handle mixed type columns to avoid Parquet type conversion errors
            for col in df.columns:
                if df[col].dtype == 'object':
                    # For object type (usually mixed type), convert uniformly to string
                    # This avoids Parquet type conversion errors
                    df[col] = df[col].astype(str)
                    # Convert 'nan' string to None
                    df[col] = df[col].replace('nan', None)
                elif pd.api.types.is_integer_dtype(df[col]):
                    # Integer type, check for NaN, convert to string if found
                    if df[col].isna().any():
                        df[col] = df[col].astype(str)
                elif pd.api.types.is_float_dtype(df[col]):
                    # Float type, check for NaN
                    if df[col].isna().any():
                        # NaN values can be handled in Parquet, but for safety, keep as is
                        pass
            
            new_file = os.path.splitext(file_path)[0] + '.parquet'
            # Use pyarrow engine
            df.to_parquet(new_file, index=False, engine='pyarrow')
            insert_formatity_task_log_info(task_uid, f'convert file {new_file} succeed')
            os.remove(file_path)
            return {
                "from": file_path,
                "to": new_file,
                "status": "success"
            }
        except Exception as e:
            print(f"convert file {file_path} error: {e}")
            insert_formatity_task_log_error(task_uid, f"convert file {file_path} error: {e}")
            return {
                "from": file_path,
                "to": None,
                "status": "failure"
            }
    else:
        return None  # Not a target file, return None


def fix_email_links_in_html(html_content: str) -> str:
    """
    修复HTML中邮箱链接的href属性错误格式
    
    Word文档中的邮箱链接可能包含Office特有的 \o 属性，mammoth转换时
    会将这些属性错误地放到href值中，导致格式如：
    href="http://email@domain.com&quot; \o &quot;http://email@domain.com"
    或
    href="mailto:email@domain.com&quot; \o &quot;mailto:email@domain.com"
    
    此函数会提取正确的邮箱地址并修复href属性。
    
    Args:
        html_content: 原始HTML内容
        
    Returns:
        修复后的HTML内容
    """
    # 模式1: 匹配 href="http://email@domain.com&quot; \o &quot;http://email@domain.com"
    # 提取第二个邮箱地址（通常是完整的）
    pattern1 = r'href="(http://[^@]+@[^"]+?)&quot;\s*\\o\s*&quot;(http://[^"]+?)"'
    
    def replace_email_link1(match):
        email1 = match.group(1)  # 第一个邮箱地址
        email2 = match.group(2)  # 第二个邮箱地址（通常是完整的）
        # 使用第二个邮箱地址，提取邮箱部分（去掉http://）
        email = email2.replace('http://', '')
        if '@' in email:
            return f'href="mailto:{email}"'
        return match.group(0)  # 如果提取失败，保持原样
    
    # 模式2: 匹配 href="mailto:email@domain.com&quot; \o &quot;mailto:email@domain.com"
    # 移除多余的 \o 部分
    pattern2 = r'href="(mailto:[^@]+@[^"]+?)&quot;\s*\\o\s*&quot;mailto:[^"]+?"'
    
    def replace_email_link2(match):
        email = match.group(1)  # mailto:email@domain.com
        return f'href="{email}"'
    
    # 模式3: 匹配 href="http://@domain.com&quot; \o &quot;http://email@domain.com"
    # 这种情况第一个是@domain，第二个是完整邮箱
    pattern3 = r'href="http://@([^"]+?)&quot;\s*\\o\s*&quot;http://([^"]+?)"'
    
    def replace_email_link3(match):
        domain = match.group(1)  # domain部分
        email = match.group(2)   # 完整邮箱地址
        if '@' in email:
            return f'href="mailto:{email}"'
        return match.group(0)
    
    # 模式4: 匹配简单的 http://email@domain.com 格式（没有 \o 部分）
    # 转换为 mailto: 格式
    pattern4 = r'href="http://([^@]+@[^/"]+?)"'
    
    def replace_email_link4(match):
        email = match.group(1)
        # 确保是邮箱格式（包含@和域名）
        if '@' in email and '.' in email.split('@')[1]:
            return f'href="mailto:{email}"'
        return match.group(0)  # 如果不是邮箱，保持原样
    
    # 按顺序应用所有修复模式
    fixed_html = re.sub(pattern1, replace_email_link1, html_content)
    fixed_html = re.sub(pattern2, replace_email_link2, fixed_html)
    fixed_html = re.sub(pattern3, replace_email_link3, fixed_html)
    fixed_html = re.sub(pattern4, replace_email_link4, fixed_html)
    
    return fixed_html


def convert_word_to_markdown(file_path: str, task_uid) -> Optional[Dict[str, str]]:
    if file_path.lower().endswith(('.docx', '.doc')):
        insert_formatity_task_log_info(task_uid, f'Source file address：{file_path}')
        try:
            with open(file_path, "rb") as docx_file:
                result = mammoth.convert_to_html(docx_file)
                html_content = result.value
            # 修复HTML中错误的邮箱链接格式
            html_content = fix_email_links_in_html(html_content)
            markdown_content = md(html_content)
            markdown_file_path = os.path.splitext(file_path)[0] + '.md'
            with open(markdown_file_path, 'w', encoding='utf-8') as md_file:
                md_file.write(markdown_content)
            insert_formatity_task_log_info(task_uid, f'convert file {markdown_file_path} succeed')
            os.remove(file_path)
            return {
                "from": file_path,
                "to": markdown_file_path,
                "status": "success"
            }
        except Exception as e:
            print(f"convert file {file_path} error: {e}")
            insert_formatity_task_log_error(task_uid, f"convert file {file_path} error: {e}")
            return {
                "from": file_path,
                "to": None,
                "status": "failure"
            }
    else:
        return None  # Not a target file, return None


def convert_ppt_to_markdown(file_path: str, task_uid) -> Optional[Dict[str, str]]:
    if file_path.lower().endswith(('.pptx', '.ppt')):
        insert_formatity_task_log_info(task_uid, f'Source file address：{file_path}')
        try:
            prs = Presentation(file_path)
            markdown_content = ""
            for i, slide in enumerate(prs.slides):
                markdown_content += f" lantern slide {i + 1}\n\n"
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text_content = shape.text.strip()
                        if text_content:
                            if len(text_content) < 50 and '\n' not in text_content:
                                markdown_content += f"## {text_content}\n\n"
                            else:
                                markdown_content += f"{text_content}\n\n"
            markdown_file_path = os.path.splitext(file_path)[0] + '.md'
            with open(markdown_file_path, 'w', encoding='utf-8') as md_file:
                md_file.write(markdown_content)
            insert_formatity_task_log_info(task_uid, f'convert file {markdown_file_path} succeed')
            os.remove(file_path)
            return {
                "from": file_path,
                "to": markdown_file_path,
                "status": "success"
            }
        except Exception as e:
            print(f"convert file {file_path} error: {e}")
            insert_formatity_task_log_error(task_uid, f"convert file {file_path} error: {e}")
            return {
                "from": file_path,
                "to": None,
                "status": "failure"
            }
    else:
        return None  # Not a target file, return None


def convert_pdf_to_markdown(file_path: str, task_uid, mineru_api_url: Optional[str] = None, mineru_backend: Optional[str] = None) -> Optional[Dict[str, str]]:
    if file_path.lower().endswith('.pdf'):
        insert_formatity_task_log_info(task_uid, f'Source file address：{file_path}')
        try:
            from mineru.cli.common import prepare_env
            from mineru.data.data_reader_writer import FileBasedDataWriter
            from mineru.backend.vlm.vlm_middle_json_mkcontent import union_make as vlm_union_make
            from mineru.utils.enum_class import MakeMode
            
            # Priority: passed parameter > environment variable > default value
            if mineru_api_url:
                server_url = mineru_api_url
            else:
                server_url = os.getenv("MINERU_API_URL", "http://111.4.242.20:30000")
            # MinerU backend: Priority: passed parameter > environment variable > default value
            if mineru_backend:
                backend = mineru_backend
            else:
                backend = os.getenv("MINERU_BACKEND", "http-client")
            
            # Record used MinerU API address and backend for debugging
            insert_formatity_task_log_info(task_uid, f'Using MinerU API server: {server_url}, backend: {backend}')
            
            pdf_file_name = Path(file_path).stem
            temp_output_dir = Path(file_path).parent / f"_temp_pdf_convert_{pdf_file_name}"
            temp_output_dir.mkdir(exist_ok=True)
            
            # Run MinerU in separate process (using subprocess)
            result_json_path = temp_output_dir / "mineru_result.json"
            script_dir = Path(__file__).parent
            mineru_worker_script = script_dir / "mineru_worker.py"
            
            python_exe = sys.executable
            cmd = [
                python_exe,
                str(mineru_worker_script),
                file_path,
                str(temp_output_dir),
                server_url,
                backend,
                str(result_json_path)
            ]
            
            process = subprocess.Popen(
                cmd,
                cwd=str(Path(__file__).parent.parent.parent)
            )
            
            process.wait()
            
            if process.returncode != 0:
                error_msg = "MinerU subprocess failed"
                if result_json_path.exists():
                    try:
                        with open(result_json_path, 'r', encoding='utf-8') as f:
                            result_data = json.load(f)
                            if not result_data.get("success", False):
                                error_msg = result_data.get("error", error_msg)
                    except:
                        pass
                insert_formatity_task_log_error(task_uid, f'MinerU subprocess failed: {error_msg}')
                raise RuntimeError(f"MinerU subprocess failed: {error_msg}")
            
            if not result_json_path.exists():
                raise FileNotFoundError(f"Result JSON file not found: {result_json_path}")
            
            with open(result_json_path, 'r', encoding='utf-8') as f:
                result_data = json.load(f)
            
            if not result_data.get("success", False):
                error_msg = result_data.get("error", "Unknown error")
                insert_formatity_task_log_error(task_uid, f'MinerU subprocess error: {error_msg}')
                raise RuntimeError(f"MinerU subprocess error: {error_msg}")
            
            middle_json = result_data["middle_json"]
            
            # Generate Markdown content
            local_image_dir, local_md_dir = prepare_env(str(temp_output_dir), pdf_file_name, "vlm")
            md_writer = FileBasedDataWriter(local_md_dir)
            
            pdf_info = middle_json["pdf_info"]
            image_dir = ""
            md_content_str = vlm_union_make(pdf_info, MakeMode.MM_MD, image_dir)
            
            markdown_filename = f"{pdf_file_name}.md"
            md_writer.write_string(markdown_filename, md_content_str)
            
            markdown_file_path = Path(local_md_dir) / markdown_filename
            final_markdown_path = os.path.splitext(file_path)[0] + '.md'
            shutil.move(str(markdown_file_path), final_markdown_path)
            
            insert_formatity_task_log_info(task_uid, f'convert file {final_markdown_path} succeed')
            os.remove(file_path)
            
            if temp_output_dir.exists():
                shutil.rmtree(temp_output_dir)
            
            return {
                "from": file_path,
                "to": final_markdown_path,
                "status": "success"
            }
        except Exception as e:
            print(f"convert file {file_path} error: {e}")
            insert_formatity_task_log_error(task_uid, f"convert file {file_path} error: {e}")
            
            temp_output_dir = Path(file_path).parent / f"_temp_pdf_convert_{Path(file_path).stem}"
            if temp_output_dir.exists():
                try:
                    shutil.rmtree(temp_output_dir)
                except:
                    pass
            
            return {
                "from": file_path,
                "to": None,
                "status": "failure"
            }
    else:
        return None  # Not a target file, return None


def search_files(folder_path: str, types: List[int]) -> Tuple[bool, List[str]]:

    type_map: Dict[int, List[str]] = {
        0: ['.ppt', '.pptx'],  # PPT
        1: ['.doc', '.docx'],  # Word
        3: ['.xls', '.xlsx'],  # Excel
        7: ['.pdf']  # PDF
    }


    target_extensions = set()
    for file_type in types:
        if file_type in type_map:
            for ext in type_map[file_type]:
                target_extensions.add(ext.lower())


    found_files: List[str] = []

    def traverse(current_path: str) -> None:

        try:

            entries = os.listdir(current_path)

            for entry in entries:
                entry_path = os.path.join(current_path, entry)

                if os.path.isdir(entry_path):

                    traverse(entry_path)
                elif os.path.isfile(entry_path):

                    file_ext = os.path.splitext(entry)[1].lower()
                    if file_ext in target_extensions:
                        found_files.append(entry_path)

        except PermissionError:
            print(f"No permission to access the folder: {current_path}")
        except Exception as e:
            print(f"Processing path {current_path} error: {str(e)}")


    traverse(folder_path)


    return bool(len(found_files) > 0)